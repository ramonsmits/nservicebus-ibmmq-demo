#!/usr/bin/env python3
"""Build docs/video/slides.pptx from the Particular presentation template.

Maps:
  Template slide 1 (TITLE)          -> our Title slide (video scene 1)
  Template slide 2 (TITLE_AND_BODY) -> our Agenda slide (video scene 2)
  Template slide 3 (Welcome)        -> deleted
  Template slide 4 (Questions?)     -> deleted
  Template slide 5 (Wrap up)        -> rebuilt as outro (video scene 7) using TITLE_AND_BODY layout
"""
from pptx import Presentation
from copy import deepcopy

SRC = "/home/ramon/Documents/particular/Particular presentation template.pptx"
DST = "/home/ramon/src/ibmmq-bridge-demo/docs/video/slides.pptx"

prs = Presentation(SRC)


def set_placeholder_text(placeholder, lines):
    """Replace placeholder text with the given list of paragraph strings.

    Paragraph-level formatting (bullets, indent) is inherited from the slide
    layout, so clearing + resetting text keeps the deck's visual style."""
    tf = placeholder.text_frame
    tf.clear()  # leaves exactly one empty paragraph
    tf.paragraphs[0].text = lines[0]
    for extra in lines[1:]:
        new_p = tf.add_paragraph()
        new_p.text = extra


# ---------- Slide 1: Title ----------
slide1 = prs.slides[0]
title = slide1.placeholders[0]
subtitle = slide1.placeholders[1]
set_placeholder_text(title, ["NServiceBus IBM MQ Transport"])
set_placeholder_text(subtitle, ["Generally available"])

# ---------- Slide 2: Agenda ----------
slide2 = prs.slides[1]
set_placeholder_text(slide2.placeholders[0], ["In the next five minutes"])
set_placeholder_text(slide2.placeholders[1], [
    "Build a subscriber from scratch",
    "See queues and topic subscriptions appear in IBM MQ",
    "Handle and replay a failure",
    "Trace it end-to-end",
])

# ---------- Slide 5: Outro ----------
# Replace slide 5's content. Its layout is SECTION_HEADER which only has a title,
# so we need TITLE_AND_BODY for bullets. Strategy: delete slide 5 and add a new
# slide from TITLE_AND_BODY layout at the same position.
#
# For deletion-and-add, we'll:
#   1. Add a new slide with TITLE_AND_BODY layout (appended to end)
#   2. Delete original slides 3, 4, 5 (Welcome, Questions?, Wrap up)
# Python-pptx places newly added slides at end; that's fine because after
# deletion of the middle slides, the new slide naturally becomes slide 3.

title_body_layout = prs.slide_layouts[2]  # TITLE_AND_BODY
outro = prs.slides.add_slide(title_body_layout)
set_placeholder_text(outro.placeholders[0], ["There's more in NServiceBus."])
set_placeholder_text(outro.placeholders[1], [
    "Bridge — Connect IBM MQ to RabbitMQ, Azure Service Bus, SQS, or Kafka",
    "Sagas — Long-running stateful workflows with persisted state",
    "Outbox — Exactly-once-effective processing with your database",
    "Message mutators — Plug in legacy formats (EBCDIC, fixed-length records)",
    "",
    "docs.particular.net",
    "nuget.org/packages/NServiceBus.Transport.IBMMQ",
    "particular.net/blog/nservicebus-transport-ibmmq",
])

# ---------- Delete template slides 3, 4, 5 ----------
# python-pptx has no public delete API; use the XML back door.
def delete_slide(prs_, slide_index):
    xml_slides = prs_.slides._sldIdLst
    slides_list = list(xml_slides)
    slide_id = slides_list[slide_index]
    # Remove relationship from the presentation part
    rId = slide_id.rId
    prs_.part.drop_rel(rId)
    xml_slides.remove(slide_id)

# Delete from highest index first so indices stay valid.
# Original deck: slides 0..4 (5 slides). After adding the outro, slides 0..5 (6).
# We want to keep indices 0 (title), 1 (agenda), 5 (our new outro).
# Delete indices 4, 3, 2 (Wrap up, Questions?, Welcome).
for idx in [4, 3, 2]:
    delete_slide(prs, idx)

prs.save(DST)
print(f"Saved {DST}")

# Verify output
prs2 = Presentation(DST)
print(f"Output deck: {len(prs2.slides)} slides")
for i, s in enumerate(prs2.slides):
    texts = []
    for shape in s.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    texts.append(p.text.strip())
    print(f"  [{i+1}] layout={s.slide_layout.name}")
    for t in texts:
        print(f"       - {t}")
